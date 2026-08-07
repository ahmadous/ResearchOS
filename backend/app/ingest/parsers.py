"""Extraction de texte depuis des fichiers importés.

Chaque type a un parseur ; un parseur manquant renvoie une erreur claire pour CE
type sans casser les autres. Les fichiers non textuels (images, vidéos) sont
reconnus mais non extraits (stockés comme pièces jointes par le service).
"""
from __future__ import annotations

import csv
import io
import os

# Extensions -> catégorie logique.
_KIND = {
    ".pdf": "pdf",
    ".docx": "word", ".doc": "word",
    ".xlsx": "excel", ".xlsm": "excel", ".xls": "excel",
    ".pptx": "powerpoint",
    ".md": "markdown", ".markdown": "markdown",
    ".tex": "latex", ".bib": "bibtex",
    ".txt": "text", ".csv": "csv", ".tsv": "csv", ".json": "text",
    ".png": "image", ".jpg": "image", ".jpeg": "image", ".gif": "image",
    ".webp": "image", ".bmp": "image", ".svg": "image", ".tiff": "image",
    ".mp4": "video", ".mov": "video", ".avi": "video", ".mkv": "video",
    ".webm": "video", ".m4v": "video",
}
IMAGE_KINDS = {"image"}
VIDEO_KINDS = {"video"}
BINARY_KINDS = IMAGE_KINDS | VIDEO_KINDS


def kind_of(filename: str) -> str:
    return _KIND.get(os.path.splitext(filename or "")[1].lower(), "unknown")


# --- Parseurs texte ---
def _pdf(data: bytes) -> str:
    try:
        import fitz  # PyMuPDF
        with fitz.open(stream=data, filetype="pdf") as doc:
            return "\n\n".join(page.get_text() for page in doc)
    except ImportError:
        from io import BytesIO
        from PyPDF2 import PdfReader
        reader = PdfReader(BytesIO(data))
        return "\n\n".join((p.extract_text() or "") for p in reader.pages)


def _word(data: bytes) -> str:
    import docx
    doc = docx.Document(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(c.text for c in row.cells))
    return "\n".join(parts)


def _excel(data: bytes) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    lines = []
    for ws in wb.worksheets:
        lines.append(f"# Feuille : {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                lines.append(" | ".join(cells))
    return "\n".join(lines)


def _powerpoint(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    out = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"# Diapositive {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                out.append(shape.text_frame.text)
    return "\n".join(out)


def _csv(data: bytes) -> str:
    text = data.decode("utf-8", errors="ignore")
    rows = list(csv.reader(io.StringIO(text)))
    return "\n".join(" | ".join(r) for r in rows)


def _plain(data: bytes) -> str:
    return data.decode("utf-8", errors="ignore")


_PARSERS = {
    "pdf": _pdf, "word": _word, "excel": _excel, "powerpoint": _powerpoint,
    "csv": _csv, "markdown": _plain, "latex": _plain, "bibtex": _plain,
    "text": _plain,
}


def extract(filename: str, data: bytes) -> dict:
    """Renvoie {kind, text, error}. `text` vide pour images/vidéos/erreurs."""
    kind = kind_of(filename)
    if kind in BINARY_KINDS or kind == "unknown":
        return {"kind": kind, "text": "", "error": None}
    parser = _PARSERS.get(kind)
    if not parser:
        return {"kind": kind, "text": "", "error": "type non supporté"}
    try:
        return {"kind": kind, "text": parser(data).strip(), "error": None}
    except ImportError as e:
        return {"kind": kind, "text": "",
                "error": f"bibliothèque manquante pour {kind}: {e.name}"}
    except Exception as e:  # fichier corrompu / illisible
        return {"kind": kind, "text": "", "error": f"extraction impossible: {e}"}
